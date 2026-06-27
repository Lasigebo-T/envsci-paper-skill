#!/usr/bin/env python3
"""
deck_build.py — envsci-slides
Build a presentation deck (.pptx) with speaker notes from a structured outline.
Templates: conference (default) / group / defense.

Outline JSON:
{
  "title": "Talk title",
  "subtitle": "Author, affiliation",
  "template": "conference",
  "slides": [
    {"heading": "Background", "bullets": ["...", "..."],
     "notes": "speaker notes", "image": "path/to/verified_fig.png"}
  ]
}

Usage:
  py deck_build.py outline.json --out talk.pptx
  py deck_build.py --selftest

Requires: python-pptx.
"""
from __future__ import annotations

import argparse
import json
import os
import sys
from pathlib import Path

TEMPLATES = {
    "conference": {"accent": "1F7A4D"},
    "group": {"accent": "0072B2"},
    "defense": {"accent": "CC79A7"},
}


def build_deck(outline: dict):
    from pptx import Presentation
    from pptx.util import Inches

    template = outline.get("template", "conference")
    if template not in TEMPLATES:
        raise ValueError(f"unknown template {template!r}; choose {list(TEMPLATES)}")

    prs = Presentation()

    # Title slide
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = outline.get("title", "Untitled")
    from pptx.dml.color import RGBColor
    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(TEMPLATES[template]["accent"])
    if len(s.placeholders) > 1:
        s.placeholders[1].text = outline.get("subtitle", "")

    # Content slides
    for sl in outline.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = sl.get("heading", "")
        tf = slide.placeholders[1].text_frame
        tf.clear()
        bullets = sl.get("bullets", []) or [""]
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = b
        img = sl.get("image")
        if img:
            if os.path.isfile(img):
                slide.shapes.add_picture(img, Inches(5.2), Inches(1.8), height=Inches(3.5))
            else:
                print(f"WARNING: image not found, skipping: {img}", file=sys.stderr)
        notes = sl.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    return prs


def run_selftest() -> int:
    import tempfile
    from pptx import Presentation

    outline = {
        "title": "Self-test talk", "subtitle": "envsci-slides",
        "template": "conference",
        "slides": [
            {"heading": "Background", "bullets": ["point a", "point b"], "notes": "intro"},
            {"heading": "Methods", "bullets": ["sampling", "QA/QC"], "notes": "explain"},
            {"heading": "Results", "bullets": ["finding 1"], "notes": "key result"},
        ],
    }
    prs = build_deck(outline)
    with tempfile.TemporaryDirectory() as d:
        out = os.path.join(d, "selftest.pptx")
        prs.save(out)
        assert os.path.isfile(out), "pptx not written"
        chk = Presentation(out)
        assert len(chk.slides) == 4, len(chk.slides)  # 1 title + 3 content
        notes_found = sum(
            1 for sl in chk.slides
            if sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip()
        )
        assert notes_found >= 3, notes_found
    print("selftest: OK")
    return 0


def main(argv=None) -> int:
    p = argparse.ArgumentParser(
        prog="deck_build.py",
        description="Build a .pptx deck with speaker notes from a structured outline.",
    )
    p.add_argument("outline", nargs="?", help="Path to outline JSON.")
    p.add_argument("--out", default="deck.pptx")
    p.add_argument("--selftest", action="store_true")
    args = p.parse_args(argv)
    if args.selftest:
        return run_selftest()
    if not args.outline:
        p.error("an outline JSON path is required (or use --selftest)")
    outline = json.loads(Path(args.outline).read_text(encoding="utf-8"))
    prs = build_deck(outline)
    prs.save(args.out)
    print(f"wrote {args.out} ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
